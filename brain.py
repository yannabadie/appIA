import requests, os, datetime
from dotenv import load_dotenv
load_dotenv()

def log_fallback(llm_from, llm_to, error):
    with open("llm_fallback.log", "a", encoding="utf-8") as logf:
        logf.write(f"{datetime.datetime.now()} | {llm_from} => {llm_to} | {str(error)}\n")

def route_llm(messages, llm="auto"):
    tried = []
    order = [llm] if llm!="auto" else ["ollama", "openai", "gemini"]
    for try_llm in order:
        tried.append(try_llm)
        try:
            if try_llm == "ollama":
                url = "http://localhost:11434/v1/chat/completions"
                payload = {"model": os.getenv("OLLAMA_MODEL", "mistral"), "messages": messages, "stream": False}
                r = requests.post(url, json=payload, timeout=60)
                result = r.json()
                return result["choices"][0]["message"]["content"]
            elif try_llm == "openai":
                import openai
                openai.api_key = os.getenv("OPENAI_API_KEY")
                completion = openai.ChatCompletion.create(model="gpt-4o", messages=messages)
                return completion.choices[0].message.content
            elif try_llm == "gemini":
                import google.generativeai as genai
                genai.configure(api_key=os.getenv("GEMINI_API_KEY"))
                model = genai.GenerativeModel('gemini-1.5-flash')
                response = model.generate_content([m['content'] for m in messages])
                return response.text
        except Exception as e:
            if len(order) > 1 and try_llm != order[-1]:
                next_llm = order[order.index(try_llm)+1]
                log_fallback(try_llm, next_llm, e)
                continue
            else:
                log_fallback(try_llm, "none", e)
                return f"[LLM ERROR] {try_llm} down ({e})"
import openai
import os
from dotenv import load_dotenv
load_dotenv()
import json
from googleapiclient.discovery import build
from google.oauth2 import service_account
from googleapiclient.http import MediaFileUpload
import requests

# HANDLERS — Google
def create_drive_folder(folder_name):
    SCOPES = ['https://www.googleapis.com/auth/drive']
    creds = service_account.Credentials.from_service_account_file(
        os.getenv("GOOGLE_APPLICATION_CREDENTIALS"), scopes=SCOPES)
    service = build('drive', 'v3', credentials=creds)
    file_metadata = {'name': folder_name, 'mimeType': 'application/vnd.google-apps.folder'}
    folder = service.files().create(body=file_metadata, fields='id').execute()
    return f"Dossier Drive créé ! ID: {folder.get('id')}"

def create_google_doc(title, text="Document généré par Jarvis AI."):
    SCOPES = ['https://www.googleapis.com/auth/documents']
    creds = service_account.Credentials.from_service_account_file(
        os.getenv("GOOGLE_APPLICATION_CREDENTIALS"), scopes=SCOPES)
    service = build('docs', 'v1', credentials=creds)
    doc = service.documents().create(body={'title': title}).execute()
    doc_id = doc.get('documentId')
    requests = [{'insertText': {'location': {'index': 1}, 'text': text}}]
    service.documents().batchUpdate(documentId=doc_id, body={'requests': requests}).execute()
    return f"Doc créé: https://docs.google.com/document/d/{doc_id}/edit"

def create_google_sheet(title):
    SCOPES = ['https://www.googleapis.com/auth/spreadsheets']
    creds = service_account.Credentials.from_service_account_file(
        os.getenv("GOOGLE_APPLICATION_CREDENTIALS"), scopes=SCOPES)
    service = build('sheets', 'v4', credentials=creds)
    spreadsheet = {'properties': {'title': title}}
    sheet = service.spreadsheets().create(body=spreadsheet, fields='spreadsheetId').execute()
    return f"Sheet créé: https://docs.google.com/spreadsheets/d/{sheet.get('spreadsheetId')}/edit"

def create_google_slide(title):
    SCOPES = ['https://www.googleapis.com/auth/presentations']
    creds = service_account.Credentials.from_service_account_file(
        os.getenv("GOOGLE_APPLICATION_CREDENTIALS"), scopes=SCOPES)
    service = build('slides', 'v1', credentials=creds)
    pres = service.presentations().create(body={'title': title}).execute()
    slide_id = pres.get('presentationId')
    return f"Slides créé: https://docs.google.com/presentation/d/{slide_id}/edit"

def tts_google(text, lang="fr-FR"):
    from google.cloud import texttospeech
    client = texttospeech.TextToSpeechClient()
    synthesis_input = texttospeech.SynthesisInput(text=text)
    voice = texttospeech.VoiceSelectionParams(
        language_code=lang, ssml_gender=texttospeech.SsmlVoiceGender.NEUTRAL
    )
    audio_config = texttospeech.AudioConfig(audio_encoding=texttospeech.AudioEncoding.MP3)
    response = client.synthesize_speech(
        input=synthesis_input, voice=voice, audio_config=audio_config
    )
    path = f"tts_{int(os.times()[4])}.mp3"
    with open(path, "wb") as out:
        out.write(response.audio_content)
    return f"Audio TTS généré : {path}"

# HANDLERS — OpenAI/Gemini/LLM
def call_openai(prompt, model="gpt-4o"):
    api_key = os.getenv("OPENAI_API_KEY")
    
    r = client = openai.OpenAI(api_key=os.getenv('OPENAI_API_KEY')); client.chat.completions.create(
        model=model,
        messages=[{"role": "user", "content": prompt}],
        max_tokens=1024,
    )
    return r.choices[0].message.content

def call_gemini(prompt):
    # Placeholder : à compléter selon ton SDK Gemini/Palm
    return "[Gemini handler ici – à brancher sur Gemini API 2.5 Pro]"

# HANDLERS — Local files
def save_to_docx(text, filename="generated.docx"):
    from docx import Document
    doc = Document()
    doc.add_paragraph(text)
    doc.save(filename)
    return f"Docx généré: {filename}"

def save_to_pptx(text, filename="generated.pptx"):
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Présentation IA"
    slide.placeholders[1].text = text
    prs.save(filename)
    return f"PPTX généré: {filename}"

# Handler — résumé contextuel
def summarize_text(text, model="gpt-4o"):
    return call_openai(f"Résume ce texte de façon concise :\n{text}", model=model)

# Méta-router LLM : analyse du prompt pour choisir le bon outil
def get_best_llm_action(prompt):
    # Utilise OpenAI pour classifier l'intention
    intent_prompt = f"""Analyse le prompt utilisateur suivant et dis-moi s'il veut : créer_dossier, créer_doc, créer_sheet, créer_slide, tts, docx, pptx, resumer, vision, ou 'ia_pure'.
    Prompt : '{prompt}'"""
    result = call_openai(intent_prompt)
    result = result.lower()
    if "dossier" in result:
        return "create_drive_folder"
    elif "doc" in result:
        return "create_google_doc"
    elif "sheet" in result or "tableur" in result:
        return "create_google_sheet"
    elif "slide" in result or "présentation" in result:
        return "create_google_slide"
    elif "tts" in result or "voix" in result:
        return "tts_google"
    elif "docx" in result:
        return "save_to_docx"
    elif "pptx" in result:
        return "save_to_pptx"
    elif "resum" in result:
        return "summarize_text"
    elif "vision" in result or "image" in result:
        return "vision"
    else:
        return "ia_pure"

def brain(prompt, history=None):
    action = get_best_llm_action(prompt)
    try:
        if action == "create_drive_folder":
            # Nom du dossier = mot clé après "nommé" ou à la fin du prompt
            name = prompt.split("nommé")[-1].strip() if "nommé" in prompt else "Nouveau dossier"
            return create_drive_folder(name)
        elif action == "create_google_doc":
            name = prompt.split("nommé")[-1].strip() if "nommé" in prompt else "Doc IA"
            text = prompt
            return create_google_doc(name, text)
        elif action == "create_google_sheet":
            name = prompt.split("nommé")[-1].strip() if "nommé" in prompt else "Sheet IA"
            return create_google_sheet(name)
        elif action == "create_google_slide":
            name = prompt.split("nommé")[-1].strip() if "nommé" in prompt else "Slides IA"
            return create_google_slide(name)
        elif action == "tts_google":
            return tts_google(prompt)
        elif action == "save_to_docx":
            return save_to_docx(prompt)
        elif action == "save_to_pptx":
            return save_to_pptx(prompt)
        elif action == "summarize_text":
            return summarize_text(prompt)
        elif action == "vision":
            return "[Vision/image : handler à intégrer ici]"
        else:
            # Mode LLM pur (chat, questions, etc.)
            return call_openai(prompt)
    except Exception as e:
        return f"Erreur (brain handler): {e}"
def ensure_gradio_history(history):
    """
    Corrige/convertit n'importe quel historique en format Gradio 4+.
    """
    out = []
    if not history:
        return []
    for h in history:
        if isinstance(h, dict) and "role" in h and "content" in h:
            out.append(h)
        elif isinstance(h, tuple) and len(h) == 2:
            if h[0] in ("user", "assistant", "system"):
                out.append({"role": h[0], "content": h[1]})
            else:
                out.append({"role": "user", "content": str(h[0])})
                out.append({"role": "assistant", "content": str(h[1])})
        elif isinstance(h, str):
            out.append({"role": "user", "content": h})
        else:
            out.append({"role": "assistant", "content": str(h)})
    return out







